from pathlib import Path

import docx
from openpyxl import load_workbook
from pptx import Presentation


def extract_text_docx(path: Path) -> str:
    document = docx.Document(path)
    return "\n".join(paragraph.text for paragraph in document.paragraphs)


def extract_text_pptx(path: Path) -> str:
    presentation = Presentation(path)
    chunks = []
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                chunks.append(shape.text)
    return "\n".join(chunks)


def extract_text_xlsx(path: Path) -> str:
    workbook = load_workbook(path, read_only=True, data_only=True)
    chunks = []
    for sheet in workbook.worksheets:
        for row in sheet.iter_rows(values_only=True):
            cells = [str(cell) for cell in row if cell is not None]
            if cells:
                chunks.append("\t".join(cells))
    workbook.close()
    return "\n".join(chunks)
